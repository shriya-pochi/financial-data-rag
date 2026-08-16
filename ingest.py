"""
ingest.py — Sales & Revenue document ingestion pipeline.

Reads PDFs, PowerPoints, and Excel workbooks and loads them into a Chroma
vector collection so query_agent.py can answer sales/revenue questions
against them with citations (file, sheet, table).

Two ways this gets used:

1. Admin / CLI bulk load into the permanent knowledge base:
       python ingest.py /path/to/folder-or-file

2. Called from app.py, per visitor upload, against a throwaway in-memory
   collection (see `ingest_file`). Nothing a visitor uploads ever touches
   the permanent collection created by (1).

Excel handling is the interesting part: spreadsheets are messy (merged
cells, phantom columns, multiple tables stacked in one sheet, subtotal
rows). We trim junk columns, auto-detect table boundaries and headers,
and keep every row with the actual data types, so a "total revenue"
question can be answered exactly instead of guessed from a text preview.
"""

import os
import re
import json
import hashlib
import chromadb
import pandas as pd
import numpy as np
from sentence_transformers import SentenceTransformer
from PyPDF2 import PdfReader
from pptx import Presentation

DB_PATH = "./vector_db"
MAIN_COLLECTION_NAME = "sales_docs"
EMBED_MODEL = "all-MiniLM-L6-v2"

SUPPORTED_EXT = {".pdf", ".ppt", ".pptx", ".xlsx", ".xls", ".csv"}

# Sheets where fewer than this many columns actually contain data get
# trimmed to just those columns before any table detection runs. Excel
# loves to allocate thousands of phantom columns from a stray format.
MIN_MEANINGFUL_COL_FILL = 1

# A table's full row set gets embedded as one document (not chunked), so
# it must stay bounded or a single retrieved item can blow past the
# model's context window. For big tables we keep the first/last rows
# (where totals usually live) plus any row that looks like a subtotal.
MAX_ROWS_PER_TABLE = 150
HEAD_ROWS = 80
TAIL_ROWS = 40

# Words that flag a chunk/table as revenue-relevant, used only to tag
# metadata for retrieval bias — never to drop data.
FINANCIAL_KEYWORDS = re.compile(
    r"\b(revenue|sales|arr|mrr|bookings|pipeline|quota|forecast|"
    r"margin|total|q[1-4]|deal|pipeline|churn|expansion|renewal|"
    r"discount|price|pricing|units|units sold|profit|cost)\b",
    re.IGNORECASE,
)

_embedder = None


def get_embedder():
    """Singleton so the model is loaded once per process."""
    global _embedder
    if _embedder is None:
        _embedder = SentenceTransformer(EMBED_MODEL)
    return _embedder


def get_persistent_client(path=DB_PATH):
    return chromadb.PersistentClient(path=path)


def get_or_create_collection(client, name=MAIN_COLLECTION_NAME):
    return client.get_or_create_collection(name)


# =========================
# HELPERS
# =========================

def file_hash(path):
    h = hashlib.sha256()
    with open(path, "rb") as f:
        h.update(f.read())
    return h.hexdigest()


def safe_read_text(path):
    try:
        return open(path, "r", encoding="utf-8", errors="ignore").read()
    except Exception:
        return ""


def make_columns_unique(cols):
    seen = {}
    new_cols = []
    for c in cols:
        c = str(c)
        if c not in seen:
            seen[c] = 0
            new_cols.append(c)
        else:
            seen[c] += 1
            new_cols.append(f"{c}_{seen[c]}")
    return new_cols


def has_financial_signal(text):
    return bool(FINANCIAL_KEYWORDS.search(text or ""))


def cap_rows(df):
    """
    Keeps a table's row count bounded before it gets embedded as a single
    document. Always keeps rows containing "total"/"subtotal" (case
    insensitive, anywhere in the row) so aggregates survive the cap, plus
    a head and tail slice for context. Returns df unchanged if it's
    already small enough.
    """
    if len(df) <= MAX_ROWS_PER_TABLE:
        return df

    row_text = df.astype(str).apply(lambda r: " ".join(r.values).lower(), axis=1)
    total_mask = row_text.str.contains(r"\btotal\b|\bsubtotal\b", regex=True, na=False)

    keep_idx = set(df.index[total_mask])
    keep_idx.update(df.index[:HEAD_ROWS])
    keep_idx.update(df.index[-TAIL_ROWS:])

    kept = df.loc[df.index.isin(keep_idx)].sort_index().copy()
    return kept


# =========================
# EXCEL: TRIM PHANTOM COLUMNS
# =========================

def trim_to_meaningful_columns(df):
    col_fill = df.notna().sum(axis=0)
    meaningful_mask = col_fill > MIN_MEANINGFUL_COL_FILL
    trimmed = df.loc[:, meaningful_mask]
    trimmed = trimmed.dropna(how="all").reset_index(drop=True)
    return trimmed


# =========================
# EXCEL TABLE DETECTION
# =========================

def detect_tables(df):
    tables = []
    start = None
    for i, row in df.iterrows():
        if row.notna().sum() > 1:
            if start is None:
                start = i
        else:
            if start is not None and i - start > 1:
                tables.append(df.iloc[start:i])
                start = None
    if start is not None:
        tables.append(df.iloc[start:])
    return tables


def detect_header(table):
    scores = []
    for _, row in table.iterrows():
        values = row.dropna()
        if len(values) == 0:
            scores.append(0)
            continue
        text_ratio = sum(isinstance(v, str) for v in values) / len(values)
        unique_ratio = len(set(values)) / len(values)
        scores.append(text_ratio + unique_ratio)
    return int(np.argmax(scores)) if scores else 0


def clean_table(table):
    header_idx = detect_header(table)
    header = table.iloc[header_idx].astype(str).str.strip()
    data = table.iloc[header_idx + 1:].copy()
    data.columns = make_columns_unique(header)
    data = data.dropna(how="all")
    data = data.replace({np.nan: None})
    return data


def excel_tables_to_text(path):
    """
    Returns (flat_text, structured) where `structured` is a list of dicts,
    one per detected table, each carrying the sheet name, table index,
    columns, and every row (not just a preview) as records — so downstream
    aggregation questions have the real numbers, not a 20-row sample.
    """
    xls = pd.ExcelFile(path)
    texts = []
    structured = []

    for sheet in xls.sheet_names:
        raw = pd.read_excel(xls, sheet_name=sheet, header=None)
        raw = raw.replace({np.nan: None})
        raw_df = trim_to_meaningful_columns(pd.DataFrame(raw))

        if raw_df.empty:
            continue

        tables = detect_tables(raw_df) or [raw_df]

        for ti, table in enumerate(tables):
            try:
                df = clean_table(table)
                if df.empty:
                    continue

                full_row_count = len(df)
                df_capped = cap_rows(df)
                truncated = full_row_count > len(df_capped)

                lines = [f"Excel sheet={sheet} table={ti} columns={list(df.columns)}"]
                if truncated:
                    lines.append(
                        f"(showing {len(df_capped)} of {full_row_count} rows: "
                        f"first {HEAD_ROWS}, last {TAIL_ROWS}, and any total/subtotal rows)"
                    )
                for _, r in df_capped.iterrows():
                    row = [f"{c}={r.get(c)}" for c in df_capped.columns]
                    lines.append(", ".join(row))
                block_text = "\n".join(lines)
                texts.append(block_text)

                structured.append({
                    "sheet": sheet,
                    "table_id": ti,
                    "columns": list(df.columns),
                    "row_count": full_row_count,
                    "truncated": truncated,
                    "rows": df_capped.astype(str).to_dict(orient="records"),
                    "has_financial_signal": has_financial_signal(block_text),
                })
            except Exception as e:
                print(f"  [warn] table parse error sheet='{sheet}' table={ti}: {e}")

    return "\n\n".join(texts), structured


# =========================
# TEXT LOADERS (PDF / PPTX / CSV)
# =========================

def load_text(path):
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pdf":
            text = "\n".join(p.extract_text() or "" for p in PdfReader(path).pages)
            return text, None

        if ext == ".csv":
            df = pd.read_csv(path)
            return df.to_string(), None

        if ext in {".xlsx", ".xls"}:
            return excel_tables_to_text(path)

        if ext in {".ppt", ".pptx"}:
            prs = Presentation(path)
            parts = []
            for slide_no, slide in enumerate(prs.slides, start=1):
                slide_lines = [f"[Slide {slide_no}]"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_lines.append(shape.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            slide_lines.append(
                                ", ".join(c.text for c in row.cells)
                            )
                parts.append("\n".join(slide_lines))
            return "\n\n".join(parts), None

        return safe_read_text(path), None

    except Exception as e:
        print("  [warn] file read error:", e)
        return "", None


# =========================
# CHUNKING
# =========================

def chunk_text(text, max_len=900):
    chunks = []
    buf = ""
    for part in text.split("\n\n"):
        if len(buf) + len(part) < max_len:
            buf += part + "\n\n"
        else:
            if buf.strip():
                chunks.append(buf.strip())
            buf = part
    if buf.strip():
        chunks.append(buf.strip())
    return chunks


# =========================
# CORE INGEST (single file -> given collection)
# =========================

def ingest_file(file_path, collection, embedder=None, project="docs", source_label=None):
    """
    Ingests a single supported file into `collection` (any Chroma
    collection — persistent main DB or an ephemeral per-session one).
    Returns the number of embeddings written, or 0 if nothing usable
    was extracted.
    """
    embedder = embedder or get_embedder()
    ext = os.path.splitext(file_path)[1].lower()
    if ext not in SUPPORTED_EXT:
        raise ValueError(f"Unsupported file type: {ext}")

    fname = source_label or os.path.basename(file_path)
    fname_noext = os.path.splitext(fname)[0]

    raw, structured = load_text(file_path)
    if not raw.strip() and not structured:
        return 0

    h = file_hash(file_path)
    items = []

    # ---- file metadata record (lets "what's in finance.xlsx" work) ----
    file_meta = {
        "file_name": fname,
        "file_name_noext": fname_noext,
        "file_type": ext,
    }
    meta_text = json.dumps(file_meta)
    items.append((
        f"file_meta::{project}/{fname}",
        meta_text,
        embedder.encode([meta_text])[0],
        {"type": "file_meta", "filename": fname, "filename_noext": fname_noext, "hash": h},
    ))

    # ---- narrative / flat text chunks ----
    if raw.strip():
        for i, c in enumerate(chunk_text(raw)):
            if not c.strip():
                continue
            items.append((
                f"text::{project}/{fname}::{i}",
                c,
                embedder.encode([c])[0],
                {
                    "type": "text",
                    "filename": fname,
                    "filename_noext": fname_noext,
                    "hash": h,
                    "has_financial_signal": has_financial_signal(c),
                },
            ))

    # ---- structured excel tables, one embedding per sheet ----
    if structured:
        sheets_seen = {}
        for entry in structured:
            sheets_seen.setdefault(entry["sheet"], []).append(entry)

        for sheet_name, sheet_tables in sheets_seen.items():
            sheet_meta = json.dumps(sheet_tables, default=str)
            safe_sheet = re.sub(r"\s+", "_", sheet_name)
            items.append((
                f"excel_struct::{project}/{fname}::{safe_sheet}",
                sheet_meta,
                embedder.encode([sheet_meta])[0],
                {
                    "type": "excel_struct",
                    "filename": fname,
                    "filename_noext": fname_noext,
                    "sheet": sheet_name,
                    "hash": h,
                    "has_financial_signal": any(
                        t.get("has_financial_signal") for t in sheet_tables
                    ),
                },
            ))

    if not items:
        return 0

    collection.add(
        ids=[i[0] for i in items],
        documents=[i[1] for i in items],
        embeddings=[i[2].tolist() for i in items],
        metadatas=[i[3] for i in items],
    )
    return len(items)


# =========================
# BULK INGEST (directory or single file -> persistent main DB)
# =========================

def ingest_path(path, collection=None, embedder=None):
    embedder = embedder or get_embedder()
    if collection is None:
        client = get_persistent_client()
        collection = get_or_create_collection(client)

    project = os.path.basename(os.path.abspath(path))
    is_file = os.path.isfile(path)
    walker = (
        [(os.path.dirname(path) or ".", [], [os.path.basename(path)])]
        if is_file
        else os.walk(path)
    )

    total = 0
    for root, _, files in walker:
        for f in files:
            ext = os.path.splitext(f)[1].lower()
            if ext not in SUPPORTED_EXT:
                continue
            full = os.path.join(root, f)
            rel = f if is_file else os.path.relpath(full, path)
            print("Processing:", rel)
            try:
                n = ingest_file(full, collection, embedder, project=project, source_label=rel)
                print(f"  -> {n} embeddings")
                total += n
            except Exception as e:
                print(f"  [error] {rel}: {e}")

    print(f"\nIndexed {total} embeddings total")
    return total


if __name__ == "__main__":
    import argparse

    p = argparse.ArgumentParser(description="Bulk-ingest sales docs into the main knowledge base.")
    p.add_argument("path", help="File or folder of PDFs/PPTX/XLSX to ingest")
    args = p.parse_args()

    ingest_path(args.path)